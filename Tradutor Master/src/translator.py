import os
from typing import Dict, Iterable, List, Optional

from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

from docx_xml_handler import export_docx_with_xml
from text_adjuster import TextAdjuster
from utils import Token


def export_translated_document(
    source_path: str,
    tokens: Iterable[Token],
    output_path: str,
    enable_size_adjustment: bool = True,
    max_length_ratio: float = 1.5,
    adjust_font_size: bool = False,
) -> Dict[str, List[str]]:
    """
    Exporta documento traduzido PRESERVANDO formatação, imagens e estrutura.
    IMPORTANTE: SEMPRE exporta como .docx, independente do formato original.

    Args:
        source_path: Caminho do arquivo original
        tokens: Lista de tokens com traduções
        output_path: Caminho para salvar o arquivo traduzido (deve terminar em .docx)
        enable_size_adjustment: Se deve ajustar tamanho do texto
        max_length_ratio: Razão máxima permitida (traduzido/original)
        adjust_font_size: Se deve ajustar tamanho de fonte (DOCX/PPTX)

    Returns:
        Dicionário com avisos gerados durante exportação
    """
    tokens_list = list(tokens)

    # Verificar extensão do arquivo original
    ext = os.path.splitext(source_path)[1].lower()

    # IMPORTANTE: Garantir que output_path sempre seja .docx
    if not output_path.lower().endswith('.docx'):
        base = os.path.splitext(output_path)[0]
        output_path = f"{base}.docx"
        print(f"⚠ Output path corrigido para: {output_path}")

    # Se for PDF, usar DOCX temporário criado durante extração
    if ext == ".pdf":
        if tokens_list and hasattr(tokens_list[0], 'source_file'):
            # Usar o DOCX temporário como source
            source_path = tokens_list[0].source_file
            ext = ".docx"
        else:
            raise ValueError("PDF precisa ser convertido para DOCX antes de exportar")

    # SEMPRE exportar como DOCX
    if ext == ".docx":
        return _export_docx(source_path, tokens_list, output_path, enable_size_adjustment, max_length_ratio)
    else:
        # Outros formatos não suportados - APENAS .docx é permitido
        raise ValueError(
            f"Apenas arquivos .docx são suportados para exportação.\n"
            f"Arquivo original: {ext}\n"
            f"Converta o arquivo para .docx antes de traduzir."
        )


def _build_translation_map(tokens: Iterable[Token]) -> Dict[str, tuple]:
    """
    Constrói mapa de traduções com texto original.

    Returns:
        Dicionário {location: (original_text, translated_text)}
    """
    translation_map: Dict[str, tuple] = {}
    for token in tokens:
        translation = token.translation

        # Proteção: garantir que translation é string
        if isinstance(translation, list):
            print(f"  AVISO: translation é lista no token {token.location}, convertendo para string")
            translation = translation[0] if translation else ""
        elif not isinstance(translation, str):
            print(f"  AVISO: translation não é string no token {token.location}: {type(translation)}")
            translation = str(translation) if translation else ""

        if translation and translation.strip():
            translation_map[token.location] = (token.text, translation)
    return translation_map


def _export_docx(
    source_path: str,
    tokens: List[Token],
    output_path: str,
    enable_size_adjustment: bool = True,
    max_length_ratio: float = 1.5,
) -> Dict[str, List[str]]:
    """
    Exporta DOCX usando manipulação XML direta para preservação RIGOROSA:
    - Formatação completa (negrito, itálico, cores, fontes)
    - Imagens e shapes
    - Quebras de página
    - Tabelas com linhas delimitadoras
    - Headers e Footers
    - Número EXATO de páginas mantido
    - TEXTO COMPLETO sem truncamento (XML preserva estrutura)
    """
    warnings_map = {}

    # IMPORTANTE: NÃO truncar textos em DOCX
    # O XML preserva a estrutura, então texto maior não quebra o layout
    # Apenas gerar avisos se texto ficar muito maior
    if enable_size_adjustment:
        for token in tokens:
            if token.translation and token.translation.strip():
                original_len = len(token.text)
                translated_len = len(token.translation)
                ratio = translated_len / original_len if original_len > 0 else 1.0

                # Apenas avisar se texto cresceu muito, MAS NÃO TRUNCAR
                if ratio > max_length_ratio:
                    warning = f"Texto traduzido {ratio:.1f}x maior que original ({original_len} -> {translated_len} chars)"
                    warnings_map[token.location] = [warning]

    # Exportar usando manipulação XML direta - TEXTO COMPLETO
    export_docx_with_xml(source_path, tokens, output_path)

    return warnings_map


def _export_pptx(
    source_path: str,
    translation_map: Dict[str, tuple],
    output_path: str,
    adjuster: Optional[TextAdjuster] = None,
) -> Dict[str, List[str]]:
    """
    Exporta PPTX substituindo APENAS o texto dos runs, preservando formatação completa.
    """
    pres = Presentation(source_path)
    warnings_map = {}

    for s_idx, slide in enumerate(pres.slides):
        for shape_idx, shape in enumerate(slide.shapes):
            if not getattr(shape, "has_text_frame", False):
                continue

            text_frame = shape.text_frame
            for para_idx, para in enumerate(text_frame.paragraphs):
                for run_idx, run in enumerate(para.runs):
                    location = f"S{s_idx}SH{shape_idx}P{para_idx}R{run_idx}"
                    trans_data = translation_map.get(location)
                    if trans_data:
                        original_text, translated_text = trans_data

                        if adjuster:
                            result = adjuster.adjust_text(original_text, translated_text)
                            translated_text = result.adjusted_text
                            if result.warnings:
                                warnings_map[location] = result.warnings

                        run.text = translated_text

    pres.save(output_path)
    return warnings_map


def _export_xlsx(
    source_path: str,
    translation_map: Dict[str, tuple],
    output_path: str,
    keep_vba: bool = False,
    adjuster: Optional[TextAdjuster] = None,
) -> Dict[str, List[str]]:
    """Exporta XLSX com ajuste de tamanho."""
    wb = load_workbook(source_path, keep_vba=keep_vba)
    warnings_map = {}

    for sheet in wb:
        for row_idx, row in enumerate(sheet.iter_rows(), start=1):
            for col_idx, cell in enumerate(row, start=1):
                location = f"{sheet.title}!R{row_idx}C{col_idx}"
                trans_data = translation_map.get(location)
                if trans_data:
                    original_text, translated_text = trans_data
                    warnings = []

                    if adjuster:
                        result = adjuster.adjust_text(original_text, translated_text)
                        translated_text = result.adjusted_text
                        warnings = result.warnings

                    cell.value = translated_text

                    if warnings:
                        warnings_map[location] = warnings

    wb.save(output_path)
    wb.close()
    return warnings_map


def _export_txt(
    source_path: str,
    translation_map: Dict[str, tuple],
    output_path: str,
    adjuster: Optional[TextAdjuster] = None,
) -> Dict[str, List[str]]:
    """Exporta TXT com ajuste de tamanho."""
    with open(source_path, "r", encoding="utf-8") as src_file:
        lines = src_file.readlines()

    new_lines = []
    warnings_map = {}

    for idx, line in enumerate(lines, start=1):
        location = f"Linha {idx}"
        trans_data = translation_map.get(location)
        if trans_data:
            original_text, translated_text = trans_data
            warnings = []

            if adjuster:
                result = adjuster.adjust_text(original_text, translated_text)
                translated_text = result.adjusted_text
                warnings = result.warnings

            ending = ""
            if line.endswith("\r\n"):
                ending = "\r\n"
            elif line.endswith("\n"):
                ending = "\n"
            new_lines.append(f"{translated_text}{ending}")

            if warnings:
                warnings_map[location] = warnings
        else:
            new_lines.append(line)

    with open(output_path, "w", encoding="utf-8") as out_file:
        out_file.writelines(new_lines)

    return warnings_map
