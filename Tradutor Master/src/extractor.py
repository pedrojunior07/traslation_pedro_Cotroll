import hashlib
import os
import logging
from typing import List, Callable, Optional

from docx import Document
from openpyxl import load_workbook
from pdf2docx import Converter
from pptx import Presentation

try:
    from .docx_xml_handler import extract_tokens_from_docx_xml
    from .utils import Token
except ImportError:
    from docx_xml_handler import extract_tokens_from_docx_xml
    from utils import Token


def extract_tokens(
    file_path: str,
    progress_callback: Optional[Callable[[str, float], None]] = None
) -> List[Token]:
    """
    Extrai tokens de um arquivo, com callback opcional de progresso.

    Args:
        file_path: Caminho do arquivo
        progress_callback: Função (mensagem, progresso_0_100) para atualizar UI
    """
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".docx":
        return _extract_docx(file_path, progress_callback)
    if ext in (".pptx", ".ppsx"):
        return _extract_pptx(file_path, progress_callback)
    if ext in (".xlsx", ".xlsm"):
        return _extract_xlsx(file_path, progress_callback)
    if ext in (".txt",):
        return _extract_txt(file_path, progress_callback)
    if ext in (".pdf",):
        return _extract_pdf(file_path, progress_callback)
    raise ValueError(f"Extensao nao suportada: {ext}")


def _extract_docx(file_path: str, progress_callback: Optional[Callable[[str, float], None]] = None) -> List[Token]:
    """
    Extrai tokens usando manipulação XML direta para preservação RIGOROSA.
    Cada elemento <w:t> no XML vira um token separado.
    Preserva ABSOLUTAMENTE TODA formatação, tabelas, imagens, quebras de página.
    """
    if progress_callback:
        progress_callback("Extraindo tokens do documento...", 50)
    return extract_tokens_from_docx_xml(file_path)


def _extract_pptx(file_path: str, progress_callback: Optional[Callable[[str, float], None]] = None) -> List[Token]:
    """
    Extrai tokens do PowerPoint preservando estrutura.
    """
    if progress_callback:
        progress_callback("Abrindo apresentação...", 10)

    pres = Presentation(file_path)
    tokens: List[Token] = []
    total_slides = len(pres.slides)

    for s_idx, slide in enumerate(pres.slides):
        if progress_callback:
            progress = 10 + (s_idx / total_slides * 80)
            progress_callback(f"Processando slide {s_idx + 1}/{total_slides}...", progress)

        for shape_idx, shape in enumerate(slide.shapes):
            if not getattr(shape, "has_text_frame", False):
                continue

            text_frame = shape.text_frame
            for para_idx, para in enumerate(text_frame.paragraphs):
                for run_idx, run in enumerate(para.runs):
                    text = run.text
                    if text and text.strip():
                        location = f"S{s_idx}SH{shape_idx}P{para_idx}R{run_idx}"
                        tokens.append(Token(file_path, location, text))

    return tokens


def _extract_xlsx(file_path: str, progress_callback: Optional[Callable[[str, float], None]] = None) -> List[Token]:
    if progress_callback:
        progress_callback("Abrindo planilha...", 10)

    wb = load_workbook(file_path, data_only=True, read_only=True)
    tokens: List[Token] = []
    total_sheets = len(wb.sheetnames)

    for sheet_idx, sheet in enumerate(wb):
        if progress_callback:
            progress = 10 + (sheet_idx / total_sheets * 80)
            progress_callback(f"Processando aba {sheet.title}...", progress)

        for row_idx, row in enumerate(sheet.iter_rows(values_only=True), start=1):
            for col_idx, value in enumerate(row, start=1):
                if isinstance(value, str):
                    text = value
                    if text and text.strip():
                        loc = f"{sheet.title}!R{row_idx}C{col_idx}"
                        tokens.append(Token(file_path, loc, text))
    wb.close()
    return tokens


def _extract_txt(file_path: str, progress_callback: Optional[Callable[[str, float], None]] = None) -> List[Token]:
    if progress_callback:
        progress_callback("Lendo arquivo de texto...", 50)

    tokens: List[Token] = []
    with open(file_path, "r", encoding="utf-8") as f:
        for idx, line in enumerate(f, start=1):
            text = line.rstrip("\r\n")
            if text and text.strip():
                tokens.append(Token(file_path, f"Linha {idx}", text))
    return tokens


def _pdf_to_docx(pdf_path: str, progress_callback: Optional[Callable[[str, float], None]] = None) -> str:
    cache_dir = os.path.join(os.path.expanduser("~"), ".tradutor_master_cache")
    os.makedirs(cache_dir, exist_ok=True)
    stamp = f"{pdf_path}|{os.path.getmtime(pdf_path)}|{os.path.getsize(pdf_path)}"
    digest = hashlib.md5(stamp.encode("utf-8")).hexdigest()
    docx_path = os.path.join(cache_dir, f"{digest}.docx")

    if os.path.exists(docx_path):
        if progress_callback:
            progress_callback("Usando conversão em cache...", 100)
        return docx_path

    if progress_callback:
        progress_callback("Iniciando conversão PDF...", 5)

    converter = Converter(pdf_path)
    try:
        if progress_callback:
            progress_callback("Convertendo PDF para DOCX...", 50)
        converter.convert(docx_path, start=0, end=None)
        if progress_callback:
            progress_callback("Conversão concluída!", 100)
    finally:
        converter.close()
    return docx_path


def _extract_pdf(file_path: str, progress_callback: Optional[Callable[[str, float], None]] = None) -> List[Token]:
    if progress_callback:
        progress_callback("Convertendo PDF para DOCX...", 10)

    docx_path = _pdf_to_docx(file_path, progress_callback)

    if progress_callback:
        progress_callback("Extraindo tokens do documento...", 80)

    tokens = _extract_docx(docx_path, progress_callback)
    for token in tokens:
        token.source_original = file_path  # PDF original
        # Guardar também o caminho do DOCX temporário para exportação
        token.source_file = docx_path
    return tokens
